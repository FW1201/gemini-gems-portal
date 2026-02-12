from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


NAVY = RGBColor(21, 43, 77)
MID_BLUE = RGBColor(35, 74, 129)
LIGHT_BG = RGBColor(245, 248, 252)
TEXT_DARK = RGBColor(35, 35, 35)


def add_top_bar(slide, title: str, subtitle: str = "") -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.9)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tf = bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(27)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Times New Roman"
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(224, 234, 247)
        p2.font.name = "Microsoft JhengHei"


def add_title_body_slide(prs: Presentation, title: str, points: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()

    add_top_bar(slide, title)

    box = slide.shapes.add_textbox(Inches(0.85), Inches(1.35), Inches(11.9), Inches(5.9))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(24)
        p.font.name = "Microsoft JhengHei"
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(14)


def add_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.65), Inches(13.333), Inches(0.85)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = MID_BLUE
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.8), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "A Corpus-Based Study of Semantic Prosody\nin ràng and shǐ Constructions"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Times New Roman"
    p.alignment = PP_ALIGN.LEFT

    sub_box = slide.shapes.add_textbox(Inches(0.85), Inches(4.2), Inches(12), Inches(1.7))
    tf2 = sub_box.text_frame
    tf2.word_wrap = True
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.text = (
        "學術論文簡報\n"
        "Chi Wu (吳奇) & Huichen S. Hsiao (蕭惠貞)\n"
        "CASLAR, 2025, 14(2): 183-200 | DOI: 10.1515/caslar-2025-2002"
    )
    p2.font.size = Pt(20)
    p2.font.name = "Microsoft JhengHei"
    p2.font.color.rgb = RGBColor(220, 230, 246)
    p2.alignment = PP_ALIGN.LEFT


def add_table_like_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()

    add_top_bar(slide, "關鍵統計結果")

    table = slide.shapes.add_table(5, 2, Inches(0.9), Inches(1.5), Inches(11.4), Inches(4.9)).table
    table.columns[0].width = Inches(5.5)
    table.columns[1].width = Inches(5.9)

    rows = [
        ("分析對象", "ràng 與 shǐ 構式，各 1,000 例"),
        ("語料來源", "COCT Written Corpus 2020"),
        ("卡方檢定", "χ²(6, N=1000) = 52.68, p < .001"),
        ("效果量", "Cramer's V = 0.23 (small-to-medium)"),
        ("重點差異", "N2→N1: ràng 15.6% > shǐ 6.2% (φ = 0.15)"),
    ]

    for r, (k, v) in enumerate(rows):
        c0 = table.cell(r, 0)
        c1 = table.cell(r, 1)
        c0.text = k
        c1.text = v
        for cell, bold in ((c0, True), (c1, False)):
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Microsoft JhengHei"
                p.font.size = Pt(19)
                p.font.bold = bold
                p.font.color.rgb = TEXT_DARK


def add_stacked_change_chart(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()
    add_top_bar(slide, "語意韻穩定度比較（Table 6）")

    chart_data = CategoryChartData()
    chart_data.categories = ["ràng", "shǐ"]
    chart_data.add_series("有變化", (34.8, 25.2))
    chart_data.add_series("無變化", (65.2, 74.8))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(1.1),
        Inches(1.5),
        Inches(8.6),
        Inches(4.8),
        chart_data,
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.maximum_scale = 100
    chart.value_axis.minimum_scale = 0
    chart.value_axis.has_major_gridlines = True

    for series in chart.series:
        series.has_data_labels = True

    note = slide.shapes.add_textbox(Inches(10.0), Inches(1.7), Inches(3.0), Inches(4.8))
    tf = note.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = (
        "解讀重點\n"
        "• shǐ 較穩定（無變化 74.8%）\n"
        "• ràng 較易產生語意韻位移\n"
        "• 支持兩構式功能分化"
    )
    p.font.name = "Microsoft JhengHei"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_DARK


def add_pattern_chart(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()
    add_top_bar(slide, "主要變化型態比較")

    chart_data = CategoryChartData()
    chart_data.categories = ["N2→N1", "P→N2", "N2→P", "N1→N2", "P→N1", "N1→P"]
    chart_data.add_series("ràng", (15.6, 3.6, 9.6, 4.2, 1.0, 0.8))
    chart_data.add_series("shǐ", (6.2, 2.8, 5.6, 4.0, 3.6, 3.0))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.9),
        Inches(1.4),
        Inches(8.9),
        Inches(4.9),
        chart_data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = True

    for series in chart.series:
        series.has_data_labels = True

    text = slide.shapes.add_textbox(Inches(10.0), Inches(1.65), Inches(3.0), Inches(4.8))
    tf = text.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = (
        "重點發現\n"
        "• 最大差異：N2→N1\n"
        "  ràng 15.6% vs shǐ 6.2%\n"
        "• ràng 更常引發\n"
        "  對受事者的負向再詮釋\n"
        "• shǐ 較偏直接、穩定致使"
    )
    p.font.name = "Microsoft JhengHei"
    p.font.size = Pt(17)
    p.font.color.rgb = TEXT_DARK


def build_presentation(output_path: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_cover(prs)
    add_title_body_slide(
        prs,
        "研究背景與問題意識",
        [
            "ràng 與 shǐ 為漢語高頻分析型致使構式，但細緻語意差異仍未充分釐清。",
            "既有研究多聚焦搭配詞與語體分布，較少檢驗語意韻（semantic prosody）的動態變化。",
            "本研究關注兩構式在「語意偏好」與「語意韻位移」上的系統差異。",
        ],
    )
    add_title_body_slide(
        prs,
        "研究問題（RQs）",
        [
            "RQ1：ràng 與 shǐ 在動詞搭配的語意偏好有何差異？",
            "RQ2：語意韻是否會從子句左側到右側產生不同方向的位移？",
            "核心假設：兩構式不僅搭配不同，且在評價色彩的變化路徑上也存在可量化差異。",
        ],
    )
    add_title_body_slide(
        prs,
        "資料來源與研究設計",
        [
            "語料：COCT Written Corpus 2020（約 1,130 萬詞）。",
            "樣本：ràng 構式 1,000 例；shǐ 構式 1,000 例（人工排除非致使用法）。",
            "流程：CQP 擷取 → collostructional analysis → semantic prosody 標註與比較。",
        ],
    )
    add_title_body_slide(
        prs,
        "分析方法",
        [
            "Simple collexeme analysis：檢驗各構式的顯著搭配詞與語意核心。",
            "Distinctive collexeme analysis：比較重疊搭配詞在兩構式中的吸引/排斥強度。",
            "Semantic prosody：依 Stubbs (1995) 將語意韻分為正向、負向、中性，檢視位移。",
        ],
    )
    add_title_body_slide(
        prs,
        "結果一：語意偏好（Semantic Preference）",
        [
            "ràng 顯著偏向認知/智力相關動詞（如 think, know, learn 類）。",
            "shǐ 顯著偏向心理狀態動詞（如 nervous, worried, upset 類）。",
            "整體上 shǐ 涵蓋更廣動詞類別，語意域較寬；ràng 偏向外顯與認知導向事件。",
        ],
    )
    add_title_body_slide(
        prs,
        "結果二：重疊搭配詞的功能分化",
        [
            "重疊 collexeme 在兩構式中呈現不同吸引強度，顯示功能分工而非單純同義替換。",
            "ràng：更常導向認知與主觀詮釋；shǐ：更常導向心理/情緒反應。",
            "此差異支持 construction grammar 對 form-meaning pairing 的細緻區分。",
        ],
    )

    add_table_like_slide(prs)
    add_stacked_change_chart(prs)
    add_pattern_chart(prs)

    add_title_body_slide(
        prs,
        "結論與理論意涵",
        [
            "兩構式皆可表達致使，但在語意偏好與語意韻演變上呈現穩定差異。",
            "ràng：變異較高、較常朝負向或再評價方向位移；shǐ：語意韻更穩定。",
            "研究補足漢語致使研究中「構式語意韻」的量化證據，強化語法與語用整合分析。",
        ],
    )
    add_title_body_slide(
        prs,
        "限制與未來研究",
        [
            "本研究聚焦書面語料，未涵蓋口語互動中的即時語用效果。",
            "後續可擴展至 jiào、lìng 等致使標記，並比較跨語體/跨地域分布。",
            "教學應用：可據此設計 ràng/shǐ 精準選用的語意與語用教學任務。",
        ],
    )

    prs.save(output_path)


if __name__ == "__main__":
    build_presentation("outputs/學術風格簡報_讓與使語意韻研究.pptx")
